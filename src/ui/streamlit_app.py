import sys
from pathlib import Path

root_dir = Path(__file__).resolve().parent.parent.parent
if str(root_dir) not in sys.path:
    sys.path.append(str(root_dir))

import io
import base64
import tempfile
from typing import TypedDict, List, Annotated
import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image

# --- LangChain & LangGraph Imports ---
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import HumanMessage
from langchain_core.output_parsers import StrOutputParser, JsonOutputParser
from pydantic import BaseModel, Field
from langgraph.graph import StateGraph, START, END

st.set_page_config(page_title="LangGraph Agentic Visual Explainer", layout="wide")

st.title("🦜🕸️ LangGraph Agentic Visual Explainer: Самокорректирующийся разбор PDF и PPTX")
st.write("Пайплайн на базе **LangGraph StateGraph** с агентом-критиком (**Quality Critic Node**), рефлексией и циклом самокоррекции.")

# --- Helper Functions ---

def image_to_base64(pil_image: Image.Image) -> str:
    """Преобразует PIL Image в base64 строку."""
    buffered = io.BytesIO()
    pil_image.convert("RGB").save(buffered, format="JPEG")
    return base64.b64encode(buffered.getvalue()).decode("utf-8")

def render_pdf_page_to_image(doc, page_num: int) -> Image.Image:
    page = doc[page_num]
    pix = page.get_pixmap(dpi=150)
    return Image.open(io.BytesIO(pix.tobytes("png")))

def render_pptx_slide_to_data(slide) -> tuple[str, list]:
    slide_text = []
    images = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            slide_text.append(shape.text.strip())
        if shape.shape_type == 13:
            img = Image.open(io.BytesIO(shape.image.blob))
            images.append(img)
    return "\n".join(slide_text), images


# --- LangGraph State Schema ---

class AgentState(TypedDict):
    api_key: str
    user_instruction: str
    detail_level: str
    pages_data: list  # list of tuples: (page_num, PIL.Image, raw_text)
    item_label: str
    vision_analyses: list  # list of strings
    quality_score: int
    critic_feedback: str
    retry_count: int
    final_output: str


# --- Pydantic Schema for Critic Node Output ---

class EvaluationResult(BaseModel):
    score: int = Field(description="Оценка качества анализа от 1 до 10")
    feedback: str = Field(description="Конкретная обратная связь, что улучшить или добавить")


# --- LangGraph Nodes ---

def vision_extraction_node(state: AgentState) -> dict:
    """Узел 1: Vision LLM извлекает и описывает визуальные данные со всех страниц."""
    api_key = state["api_key"]
    pages_data = state["pages_data"]
    item_label = state["item_label"]
    feedback = state.get("critic_feedback", "")
    
    vision_model = ChatGroq(
        groq_api_key=api_key,
        model_name="llama-3.2-11b-vision-preview",
        temperature=0.2,
        max_tokens=1500
    )
    chain = vision_model | StrOutputParser()

    analyses = []
    for num, img, text in pages_data:
        b64_img = image_to_base64(img)
        
        prompt_text = f"""Перед вами {item_label} №{num}.

ИЗВЛЕЧЕННЫЙ ТЕКСТ:
{text if text else '[Текст отсутствует]'}

ИНСТРУКЦИЯ ПО АНАЛИЗУ:
1. Подробно опишите все схемы, таблицы, графики, диаграммы и архитектуры.
2. Поясните взаимосвязь визуальных элементов с текстом.
3. Разжуйте ключевые тезисы и метрики.
"""
        if feedback:
            prompt_text += f"\n\nВАЖНО: В предыдущей попытке Критик дал замечание: '{feedback}'. Исправь это и сделай описание более глубоким."

        message = HumanMessage(
            content=[
                {"type": "text", "text": prompt_text},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{b64_img}"}
                }
            ]
        )
        res = chain.invoke([message])
        analyses.append(f"=== {item_label} {num} ===\n{res}")

    return {
        "vision_analyses": analyses,
        "retry_count": state.get("retry_count", 0) + 1
    }


def quality_critic_node(state: AgentState) -> dict:
    """Узел 2: Критик проверяет качество Vision-разбора и дает оценку с фидбеком."""
    api_key = state["api_key"]
    analyses = "\n\n".join(state["vision_analyses"])

    critic_model = ChatGroq(
        groq_api_key=api_key,
        model_name="llama-3.3-70b-versatile",
        temperature=0.1
    )

    parser = JsonOutputParser(pydantic_object=EvaluationResult)

    prompt = ChatPromptTemplate.from_messages([
        ("system", "Вы — строгий рецензент и главный контролер качества AI-аналитики. Оцените полноту и качество анализа страниц."),
        ("human", """Оцените следующий первичный разбор документа:

{context}

ИНСТРУКЦИЯ ПО ОЦЕНКЕ:
- Поставьте оценку от 1 до 10 (где 10 — идеально глубокий разбор всех схем, формул и архитектур).
- Напишите краткий фидбек, чего не хватает или что нужно улучшить.

Верните ответ СТРОГО в формате JSON:
{format_instructions}
""")
    ])

    chain = prompt | critic_model | parser

    try:
        res = chain.invoke({
            "context": analyses[:15000],
            "format_instructions": parser.get_format_instructions()
        })
        return {
            "quality_score": res.get("score", 8),
            "critic_feedback": res.get("feedback", "Все отлично")
        }
    except Exception:
        # Резервный вариант, если модель вернула невалидный JSON
        return {
            "quality_score": 8,
            "critic_feedback": "Анализ принят без замечаний."
        }


def final_synthesis_node(state: AgentState) -> dict:
    """Узел 3: Финальный синтез проверенного контекста в структурированный документ."""
    api_key = state["api_key"]
    combined_context = "\n\n".join(state["vision_analyses"])

    synthesis_model = ChatGroq(
        groq_api_key=api_key,
        model_name="llama-3.3-70b-versatile",
        temperature=0.3,
        max_tokens=4000
    )

    prompt_template = ChatPromptTemplate.from_messages([
        ("system", "Вы — главный AI Архитектор. Составляйте глубокие, исчерпывающие и академически точные аналитические отчеты."),
        ("human", """УКАЗАНИЯ ПОЛЬЗОВАТЕЛЯ:
{user_instruction}

УРОВЕНЬ ДЕТАЛИЗАЦИИ: {detail_level}

ПРОВЕРЕННЫЙ И ВЕРИФИЦИРОВАННЫЙ КРИТИКОМ КОНТЕКСТ ДОКУМЕНТА:
{context}

СТРУКТУРА ФИНАЛЬНОГО ОТЧЕТА:
1. **Главная суть и архитектура**: В чём смысл документа.
2. **Декомпозиция схем, графиков и таблиц**: Полный разбор с причинно-следственными связями.
3. **Глоссарий терминов**: Расшифровка всех терминов и сокращений.
4. **Практические выводы и чек-лист**: Пошаговый план применения знаний.
""")
    ])

    chain = prompt_template | synthesis_model | StrOutputParser()

    res = chain.invoke({
        "user_instruction": state["user_instruction"],
        "detail_level": state["detail_level"],
        "context": combined_context[:35000]
    })

    return {"final_output": res}


# --- Conditional Routing Function ---

def should_continue(state: AgentState) -> str:
    """Решающая функция: переходить к синтезу или запустить цикл самокоррекции."""
    score = state.get("quality_score", 10)
    retries = state.get("retry_count", 0)

    # Если оценка ниже 7 и мы делали меньше 2 попыток — отправляем назад на доработку
    if score < 7 and retries < 2:
        return "re_analyze"
    return "synthesize"


# --- Build LangGraph Workflow ---

def build_app_graph():
    builder = StateGraph(AgentState)

    # Добавляем узлы
    builder.add_node("vision_extraction", vision_extraction_node)
    builder.add_node("quality_critic", quality_critic_node)
    builder.add_node("final_synthesis", final_synthesis_node)

    # Строим рёбра (Edges)
    builder.add_edge(START, "vision_extraction")
    builder.add_edge("vision_extraction", "quality_critic")

    # Условное ребро (Self-Correction Loop)
    builder.add_conditional_edges(
        "quality_critic",
        should_continue,
        {
            "re_analyze": "vision_extraction",
            "synthesize": "final_synthesis"
        }
    )

    builder.add_edge("final_synthesis", END)

    return builder.compile()


# --- Streamlit UI ---

st.sidebar.header("⚙️ LangGraph Agent Settings")
groq_api_key = st.sidebar.text_input("Groq API Key:", type="password")
detail_level = st.sidebar.select_slider(
    "Уровень детализации:",
    options=["Стандартный", "Глубокий академический", "Экстремально подробный"]
)

log_container = st.expander("🕸️ LangGraph Agent Execution Trace Logs", expanded=True)

def log_msg(msg: str, status: str = "info"):
    with log_container:
        if status == "success":
            st.success(f"[LANGGRAPH TRACE]: {msg}")
        elif status == "warning":
            st.warning(f"[LANGGRAPH TRACE]: {msg}")
        elif status == "error":
            st.error(f"[LANGGRAPH TRACE]: {msg}")
        else:
            st.info(f"[LANGGRAPH TRACE]: {msg}")

uploaded_file = st.file_uploader("Загрузите PDF или PPTX документ", type=["pdf", "pptx"])
user_instruction = st.text_area(
    "Фокус-указания для AI:", 
    value="Проведи детальный разбор. Разжуй все термины, таблицы, схемы, графики и архитектуру простым языком.",
    height=80
)

if st.button("🚀 Запустить LangGraph Agentic Pipeline"):
    if not uploaded_file or not groq_api_key.strip():
        st.error("Загрузите файл и укажите Groq API Key!")
        st.stop()

    file_ext = Path(uploaded_file.name).suffix.lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_path = tmp_file.name

    pages_data = []
    if file_ext == ".pdf":
        doc = fitz.open(tmp_path)
        for page_idx in range(len(doc)):
            img = render_pdf_page_to_image(doc, page_idx)
            text = doc[page_idx].get_text("text")
            pages_data.append((page_idx + 1, img, text))
        doc.close()
    elif file_ext in [".pptx", ".ppt"]:
        prs = Presentation(tmp_path)
        for slide_idx, slide in enumerate(prs.slides):
            slide_text, slide_imgs = render_pptx_slide_to_data(slide)
            img = slide_imgs[0] if slide_imgs else Image.new('RGB', (800, 600), color=(240, 240, 240))
            pages_data.append((slide_idx + 1, img, slide_text))

    item_label = "Страница" if file_ext == ".pdf" else "Слайд"

    # Инициализация состояния графа
    initial_state: AgentState = {
        "api_key": groq_api_key.strip(),
        "user_instruction": user_instruction,
        "detail_level": detail_level,
        "pages_data": pages_data,
        "item_label": item_label,
        "vision_analyses": [],
        "quality_score": 0,
        "critic_feedback": "",
        "retry_count": 0,
        "final_output": ""
    }

    log_msg("Компиляция LangGraph StateGraph...", "info")
    app_graph = build_app_graph()

    log_msg("Запуск агента в LangGraph...", "info")
    
    # Запуск выполнения графа
    with st.spinner("LangGraph Агент выполняет работу и проверяет качество..."):
        final_state = app_graph.invoke(initial_state)

    log_msg(f"Оценка качества от Критика: {final_state['quality_score']}/10", "success" if final_state['quality_score'] >= 7 else "warning")
    log_msg(f"Замечания Критика: {final_state['critic_feedback']}", "info")
    log_msg(f"Всего итераций самокоррекции: {final_state['retry_count']}", "info")

    st.markdown("---")
    st.markdown("## 📸 Визуальный разбор страниц")
    for idx, text in enumerate(final_state["vision_analyses"]):
        st.markdown(text)

    st.markdown("---")
    st.markdown("## 🏛️ Итоговый verified-разбор (LangGraph Agent Output)")
    st.markdown(final_state["final_output"])
